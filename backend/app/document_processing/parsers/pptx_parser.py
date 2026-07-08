"""
DealOS AI — PPTX Parser.

Extracts text from PowerPoint presentations using python-pptx.
Each slide becomes a separate page for citation tracking.
"""

import io

import structlog
from pptx import Presentation

logger = structlog.get_logger(__name__)


class PptxParser:
    """Extract text content from PowerPoint presentations."""

    @staticmethod
    async def parse(content: bytes) -> list[dict]:
        """
        Parse a PPTX file and extract text by slide.

        Each slide maps to a page for citation tracking.
        """
        pages = []

        try:
            prs = Presentation(io.BytesIO(content))

            for slide_num, slide in enumerate(prs.slides, 1):
                texts = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)

                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join(
                                cell.text.strip()
                                for cell in row.cells
                                if cell.text.strip()
                            )
                            if row_text:
                                texts.append(row_text)

                if texts:
                    pages.append({
                        "page_number": slide_num,
                        "text": "\n".join(texts),
                        "metadata": {
                            "source_type": "pptx",
                            "slide_number": slide_num,
                        },
                    })

            logger.info("pptx_parsed", total_slides=len(prs.slides), slides_with_text=len(pages))

        except Exception as e:
            logger.error("pptx_parse_error", error=str(e))
            raise

        return pages
